"""
VoiceHire Pitch Deck Content Extractor
Extracts text content from the old PowerPoint pitch deck.
"""

from pptx import Presentation
import json
import os

# Path to the old pitch deck
PPTX_PATH = "pitch deck/VoiceHire_Pitch_Deck_-_Workspace_Upgrade.pptx"

def extract_pitch_deck_content(pptx_path: str) -> list[dict]:
    """Extract all slide content from PowerPoint file."""
    print(f"[INFO] Opening: {pptx_path}")
    prs = Presentation(pptx_path)

    slides_content = []

    for i, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": i,
            "title": "",
            "content": [],
            "notes": ""
        }

        # Extract shapes (text boxes, titles, content)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()

                # Try to identify title (usually first non-empty text or has title placeholder)
                if not slide_data["title"] and (
                    hasattr(shape, "is_placeholder") or
                    len(text) < 100  # Titles are usually short
                ):
                    slide_data["title"] = text
                else:
                    slide_data["content"].append(text)

        # Extract speaker notes if available
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text.strip():
                slide_data["notes"] = notes_frame.text.strip()

        slides_content.append(slide_data)
        print(f"  [OK] Slide {i}: {slide_data['title'][:50] if slide_data['title'] else '(No title)'}")

    return slides_content


def save_as_json(slides_content: list[dict], output_path: str):
    """Save extracted content as JSON."""
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(slides_content, f, indent=2, ensure_ascii=False)
    print(f"[SAVED] JSON saved: {output_path}")


def save_as_markdown(slides_content: list[dict], output_path: str):
    """Save extracted content as readable Markdown."""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("# VoiceHire Pitch Deck — Old Version (Extracted Content)\n\n")
        f.write(f"**Total Slides:** {len(slides_content)}\n\n")
        f.write("---\n\n")

        for slide in slides_content:
            # Slide header
            title = slide['title'] if slide['title'] else f"Slide {slide['slide_number']}"
            f.write(f"## Slide {slide['slide_number']}: {title}\n\n")

            # Slide content
            if slide['content']:
                for content_item in slide['content']:
                    # Check if it's a bullet list (contains line breaks or bullet points)
                    if '\n' in content_item or '•' in content_item or '-' in content_item:
                        f.write(f"{content_item}\n\n")
                    else:
                        f.write(f"{content_item}\n\n")

            # Speaker notes
            if slide['notes']:
                f.write(f"**Speaker Notes:**\n{slide['notes']}\n\n")

            f.write("---\n\n")

    print(f"[SAVED] Markdown saved: {output_path}")


def main():
    print("=" * 60)
    print("  VoiceHire Pitch Deck Content Extractor")
    print("=" * 60)
    print()

    # Check if file exists
    if not os.path.exists(PPTX_PATH):
        print(f"[ERROR] File not found: {PPTX_PATH}")
        print("Please ensure the old pitch deck file exists at this path.")
        return

    # Extract content
    slides_content = extract_pitch_deck_content(PPTX_PATH)

    # Save outputs
    json_path = "old_pitch_deck_content.json"
    md_path = "old_pitch_deck_content.md"

    save_as_json(slides_content, json_path)
    save_as_markdown(slides_content, md_path)

    # Summary
    print()
    print("=" * 60)
    print("  EXTRACTION COMPLETE")
    print("=" * 60)
    print(f"[OK] Extracted {len(slides_content)} slides")
    print(f"[JSON] JSON output: {json_path}")
    print(f"[SAVED] Markdown output: {md_path}")
    print()
    print("Next steps:")
    print("1. Review old_pitch_deck_content.md for readable format")
    print("2. Use old_pitch_deck_content.json for structured parsing")
    print("3. Compare against current implementation features")
    print("=" * 60)


if __name__ == "__main__":
    main()
